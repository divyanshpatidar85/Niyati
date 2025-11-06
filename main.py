import os
import io
import time
import textwrap
from typing import List, Tuple
 
import streamlit as st
import pandas as pd
 
# File-type libs
import pdfplumber
import docx2txt
from pptx import Presentation
 
# LangChain / LLM
from langchain_groq import ChatGroq
from langchain.prompts import PromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain.chains import RetrievalQA
 
# Vector DB (Chroma)
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings
 
# -----------------------------
# Streamlit Page Setup
# -----------------------------
st.set_page_config(page_title="Business Report Summarizer", page_icon="📊", layout="wide")
st.title("📊 Business Report Summarizer")
st.caption("Upload any report (TXT, DOCX, PDF, Excel, PPTX) — get an instant executive summary. \
           Works with business reports, dashboards, research papers, or even raw tables.")
 
# -----------------------------
# Sidebar: API Key + Options
# -----------------------------
with st.sidebar:
    st.header("⚙️ Settings")
    groq_api_key = st.text_input("Groq API Key", type="password", help="Get from console.groq.com → API Keys")
    model_name = st.selectbox(
        "Groq Model",
        options=[
            "llama-3.1-8b-instant",    # fast, budget-friendly
            "llama-3.3-70b-versatile", # higher quality
            "gemma2-9b-it",            # Google Gemma replacement
        ],
        index=0,
    )
    style = st.selectbox(
        "Summary Style",
        ["Executive (short paragraphs)", "Bullet Points (crisp)", "Hybrid (bullets + brief context)"]
    )
    max_words = st.slider("Target Summary Length (words)", min_value=120, max_value=600, value=300, step=30)
 
# -----------------------------
# File Uploader
# -----------------------------
SUPPORTED_TYPES = ["txt", "docx", "pdf", "xlsx", "xls", "csv", "pptx"]
uploaded = st.file_uploader(
    "Upload your Business Report",
    type=SUPPORTED_TYPES,
)
 
# -----------------------------
# Extraction Helpers
# -----------------------------
def extract_from_txt(file) -> str:
    try:
        raw = file.read()
        for enc in ("utf-8", "utf-16", "latin-1"):
            try:
                return raw.decode(enc)
            except Exception:
                continue
        return raw.decode(errors="ignore")
    except Exception as e:
        return f"[Extraction Error - TXT] {e}"
 
def extract_from_docx(file) -> str:
    try:
        data = file.read()
        bio = io.BytesIO(data)
        text = docx2txt.process(bio)
        return text or ""
    except Exception as e:
        return f"[Extraction Error - DOCX] {e}"
 
def extract_from_pdf(file) -> str:
    try:
        text_parts = []
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                t = page.extract_text() or ""
                text_parts.append(t)
        text = "\n".join(text_parts)
        if not text.strip():
            return "⚠️ No selectable text found. (Might be scanned PDF)"
        return text
    except Exception as e:
        return f"[Extraction Error - PDF] {e}"
 
def extract_from_excel(file, ext: str) -> Tuple[str, List[pd.DataFrame]]:
    try:
        dfs: List[pd.DataFrame] = []
        if ext == "csv":
            df = pd.read_csv(file)
            dfs.append(df)
        else:
            xls = pd.ExcelFile(file)
            for sheet in xls.sheet_names:
                dfs.append(xls.parse(sheet))
        text_blocks = []
        for idx, df in enumerate(dfs, start=1):
            head_str = df.head(10).to_string(index=False)
            try:
                desc = df.describe(include=["number"]).to_string()
            except Exception:
                desc = "(No numeric summary available)"
            text_blocks.append(
                f"[Sheet #{idx}]\nColumns: {list(df.columns)}\nSample:\n{head_str}\n\nNumeric Summary:\n{desc}\n"
            )
        long_text = "\n\n".join(text_blocks)
        return long_text, dfs
    except Exception as e:
        return f"[Extraction Error - Excel/CSV] {e}", []
 
def extract_from_pptx(file) -> str:
    try:
        prs = Presentation(file)
        buffer = []
        for sidx, slide in enumerate(prs.slides, start=1):
            buffer.append(f"\n--- Slide {sidx} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    buffer.append(shape.text)
        return "\n".join(buffer)
    except Exception as e:
        return f"[Extraction Error - PPTX] {e}"
 
def build_text_from_upload(uploaded_file) -> Tuple[str, List[pd.DataFrame]]:
    ext = uploaded_file.name.split(".")[-1].lower()
    dfs: List[pd.DataFrame] = []
 
    if ext == "txt":
        text = extract_from_txt(uploaded_file)
    elif ext == "docx":
        text = extract_from_docx(uploaded_file)
    elif ext == "pdf":
        text = extract_from_pdf(uploaded_file)
    elif ext in ("xlsx", "xls", "csv"):
        text, dfs = extract_from_excel(uploaded_file, ext)
    elif ext == "pptx":
        text = extract_from_pptx(uploaded_file)
    else:
        text = "❌ Unsupported file type."
 
    if not isinstance(text, str):
        text = str(text)
    return text, dfs
 
# -----------------------------
# Summarization with VectorDB
# -----------------------------
def summarize_with_vectordb(text: str, api_key: str, model: str, style: str, max_words: int) -> str:
    if not text.strip():
        return "⚠️ No text extracted."
 
    llm = ChatGroq(model=model, groq_api_key=api_key)
 
    # Fixed chunk settings
    chunk_size = 4000
    chunk_overlap = 200
 
    # Embeddings + Chroma
    embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = splitter.split_text(text)
    vectordb = Chroma.from_texts(chunks, embeddings)
 
    retriever = vectordb.as_retriever(search_kwargs={"k": 6})
 
    # Universal Prompt
    instructions = {
        "Executive (short paragraphs)": (
            "You are an expert summarizer. Read the document carefully and write a structured summary.\n"
            "- Business report → KPIs, trends, risks, recommendations.\n"
            "- Dashboard/project spec → goals, visuals, slicers, metrics.\n"
            "- Research/academic → objectives, methods, findings, conclusions.\n"
            "- Numeric tables → highlight key patterns, trends, anomalies.\n"
            "- If unclear, scanned, or empty → note limitation.\n"
            f"Keep it around {max_words} words. Write in short paragraphs."
        ),
        "Bullet Points (crisp)": (
            "You are an expert summarizer. Create concise bullet points.\n"
            "- Business report → KPIs, trends, risks, actions.\n"
            "- Dashboard/project spec → visuals, slicers, metrics.\n"
            "- Research/academic → objectives, findings.\n"
            "- Tables → insights and anomalies.\n"
            "- If unclear → state limitation.\n"
            f"Limit ~{max_words} words."
        ),
        "Hybrid (bullets + brief context)": (
            "You are an expert summarizer. Create a hybrid summary.\n"
            "- Start with 2–3 line context.\n"
            "- Then list 6–10 bullet points with key insights.\n"
            "- Adjust based on content type: business, dashboard, research, tables.\n"
            "- If unclear → state limitation.\n"
            f"Keep it within {max_words} words."
        ),
    }[style]
 
    query = f"Summarize the following document:\n\n{text}\n\nInstructions:\n{instructions}"
 
    qa = RetrievalQA.from_chain_type(llm=llm, retriever=retriever, chain_type="stuff")
    return qa.run(query)
 
# -----------------------------
# MAIN UI FLOW
# -----------------------------
col_left, col_right = st.columns([1, 1])
 
with col_left:
    st.subheader("1) Upload & Extract")
    if uploaded is None:
        st.info("Upload a report to begin.")
    else:
        with st.spinner("Extracting content..."):
            extracted_text, dataframes = build_text_from_upload(uploaded)
 
        if dataframes:
            st.write(f"Found **{len(dataframes)}** table(s).")
            for i, df in enumerate(dataframes[:3], start=1):
                st.write(f"Table Preview #{i}")
                st.dataframe(df.head(20))
 
        st.text_area("Extracted Content (preview)", value=extracted_text[:8000], height=300)
 
with col_right:
    st.subheader("2) Summarize")
    summarize_btn = st.button("🔎 Generate Summary", use_container_width=True)
 
    if summarize_btn:
        if not groq_api_key.strip():
            st.error("Please enter your Groq API key.")
        else:
            with st.spinner("Summarizing with VectorDB…"):
                summary = summarize_with_vectordb(
                    extracted_text,
                    api_key=groq_api_key,
                    model=model_name,
                    style=style,
                    max_words=max_words,
                )
 
            st.success("Summary ready!")
            st.subheader("📝 Executive Summary")
            st.write(summary)